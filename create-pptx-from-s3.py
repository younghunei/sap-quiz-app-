import os
import json
import boto3
import io
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from datetime import datetime
from pptx.chart.data import CategoryChartData
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from copy import deepcopy

# Get S3 bucket name from environment variables
JSON_S3_BUCKET_NAME = os.getenv("JSON_S3_BUCKET_NAME")
PPTX_S3_BUCKET_NAME = os.getenv("PPTX_S3_BUCKET_NAME")
S3_PPTX = os.getenv("PPTX_TEMP")

def set_font(paragraph, name="Noto Sans KR", size=14, bold=False, italic=None, color=(0, 0, 0)):
    font = paragraph.font
    font.name = name
    font.size = Pt(size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)


def set_cell (cell, text, size=8, fore_color=(0,0,0), bold=False, font_color=(255,255,255), font_name="Noto Sans KR", column = 0 ) :
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*fore_color)
    if column == 1:
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER   # 가로(수평) 가운데 정렬
        cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 세로(수직) 가운데 정렬

    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*font_color)
        p.font.name = font_name
        p.font.bold = bold


def set_cell_border(cell, border_color="000000", border_width='12700'):
    """
    cell: pptx.table._Cell
    border_color: RGB hex string (예: "000000" for black)
    border_width: EMUs 단위 (12700 = 약 1pt)
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for line_dir in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        # 기존 선 제거
        for e in tcPr.findall(qn(line_dir)):
            tcPr.remove(e)
        # 새 선 추가
        ln = OxmlElement(line_dir)
        ln.set('w', border_width)
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', border_color)
        solidFill.append(srgbClr)
        ln.append(solidFill)
        prstDash = OxmlElement('a:prstDash')
        prstDash.set('val', 'solid')
        ln.append(prstDash)
        tcPr.append(ln)

def create_ppt_from_template(template_bucket, template_key, json_data, output_path):
    """
    S3 버킷에서 PPT 템플릿을 가져와 JSON 데이터를 기반으로 새 슬라이드를 추가한 PPT 생성
    
    Args:
        template_bucket (str): S3 버킷 이름
        template_key (str): 템플릿 PPT 객체 키
        json_data (dict): 새 슬라이드에 추가할 데이터
        output_path (str): 생성된 PPT를 저장할 경로
    """
    # S3 클라이언트 생성
    s3_client = boto3.client('s3')
    
    # S3 버킷에서 템플릿 파일 다운로드
    response = s3_client.get_object(Bucket=template_bucket, Key=template_key)
    template_content = response['Body'].read()
    
    # 템플릿을 메모리에 로드
    prs = Presentation(io.BytesIO(template_content))
    
    # 템플릿의 슬라이드 목록 복사
    slides = list(prs.slides)
    
    # 새 프레젠테이션 생성
    new_prs = Presentation()
    
    # 템플릿의 슬라이드 마스터와 레이아웃 복사
    for layout in prs.slide_masters[0].slide_layouts:
        new_prs.slide_masters[0].slide_layouts.add_slide_layout(layout)
    
    # 유지할 슬라이드 인덱스 (1번, 2번, 6번~11번)
    # 참고: 인덱스는 0부터 시작하므로 실제 슬라이드 번호에서 1을 뺀 값
    keep_slides = [0, 1] + list(range(5, 11))
    
    # 사용자 지정 슬라이드 (예: 3번 슬라이드를 사용자 지정 슬라이드로 가정)
    custom_slide_idx = 2  # 실제 템플릿의 사용자 지정 슬라이드 인덱스로 수정 필요
    
    # 1번과 2번 슬라이드 복사
    for i in [0, 1]:
        slide_layout = new_prs.slide_layouts[slides[i].slide_layout.index]
        new_slide = new_prs.slides.add_slide(slide_layout)
        
        # 슬라이드 내용 복사 (도형, 텍스트 등)
        for shape in slides[i].shapes:
            if hasattr(shape, "text"):
                for shape_new in new_slide.shapes:
                    if shape_new.name == shape.name:
                        shape_new.text = shape.text
    
    # JSON 데이터를 기반으로 새 슬라이드 생성 및 2번 슬라이드 뒤에 삽입
    custom_slide_layout = new_prs.slide_layouts[slides[custom_slide_idx].slide_layout.index]
    
    # JSON 데이터에서 슬라이드 정보 추출 및 새 슬라이드 생성
    for item in json_data["items"]:
        new_slide = new_prs.slides.add_slide(custom_slide_layout)
        
        # 슬라이드 내 도형에 데이터 채우기
        for shape in new_slide.shapes:
            if shape.name == "title" and hasattr(shape, "text"):
                shape.text = item.get("title", "")
            elif shape.name == "content" and hasattr(shape, "text"):
                shape.text = item.get("content", "")
            # 필요에 따라 더 많은 도형 처리 추가
    
    # 6번~11번 슬라이드 복사
    for i in range(5, 11):
        slide_layout = new_prs.slide_layouts[slides[i].slide_layout.index]
        new_slide = new_prs.slides.add_slide(slide_layout)
        
        # 슬라이드 내용 복사 (도형, 텍스트 등)
        for shape in slides[i].shapes:
            if hasattr(shape, "text"):
                for shape_new in new_slide.shapes:
                    if shape_new.name == shape.name:
                        shape_new.text = shape.text
    
    # 새 프레젠테이션 저장
    new_prs.save(output_path)
    
    return output_path

def lambda_handler(event, context):

    company = event["company"]
    key = event["s3key"]
    s3 = boto3.client("s3")

    s3_response = s3.get_object(Bucket=JSON_S3_BUCKET_NAME, Key=key)
    json_str = s3_response["Body"].read().decode("utf-8")
    json_data = json.loads(json_str)  # 여기까지 하면 json_data 는 list 형식

    
    account_ids = []
    service_names = []
    months = []

    for account in json_data:
        for ticket in account['tickets']:
            account_ids.append(account['account_id'])
            service_names.append(ticket['service_name'])
            month = ticket['date'][:7]  # 'YYYY-MM' 형식 추출
            months.append(month)

    from collections import Counter
    service_name_counts = Counter(service_names)
    months_counts = Counter(months)
    account_ids_counts = Counter(account_ids)


    pptx_stream = io.BytesIO()
    s3.download_fileobj(PPTX_S3_BUCKET_NAME, S3_PPTX, Fileobj=pptx_stream)

    # # 스트림 포인터를 처음으로 되돌림
    pptx_stream.seek(0)
    prs = Presentation(pptx_stream)

    # 고객사 이름
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        print(f"Slide {i+1} uses layout: {layout_name}")

        if i == 0 :
            # 첫 번째 장
            slide.shapes[0].text = f"2025년 1분기 기술지원 보고서\n - {company} \nCloud Support Service Group\n\n\n\n\nTechnical Support Service 1 Team"
            title,customer,group, _, _, _, _,team = slide.shapes[0].text_frame.paragraphs
            set_font(title, size=33, color=(255, 255, 255))
            set_font(customer, size=20, color=(255, 255, 255))
            set_font(group, size=28, color=(255, 153, 0))
            set_font(team, size=15, color=(255, 255, 255))

        elif i == 1 :
            # 둘 번째 장 목차
            print("slide 1 or 2 or 3")
        elif i == 5 :
            print("마지막 장은 그냥 패스")
            # 마지막 장
        elif i == 2 :
            content_slide_layout = prs.slide_layouts[13]
            slide = prs.slides.add_slide(content_slide_layout)

            title = slide.placeholders[0]
            title_tf = title.text_frame
            title_tf.text = f"1. Number of Tickets Graph"
            set_font(title_tf.paragraphs[0], size=18, bold=True)

            months_sorted = sorted(months_counts.keys())
            months_counts_sorted = [months_counts[m] for m in months_sorted]

            # pptx 내장 차트 데이터 생성
            chart_data = CategoryChartData()
            chart_data.categories = months_sorted
            chart_data.add_series('티켓 건수', months_counts_sorted)

            # 차트 추가 (위치 및 크기는 필요에 따라 조정)
            x, y, cx, cy = Inches(0.5), Inches(1), Inches(8), Inches(4.2)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            # 차트 제목/축 설정 등은 필요에 따라 추가
            # chart.has_legend = True
            

            chart.has_title = True
            chart.chart_title.text_frame.text = "월별 티켓 건수"

            chart.category_axis.has_major_gridlines = False
            # category_axis = chart.category_axis
            # category_axis.has_title = True
            # category_axis.axis_title.text_frame.text = "월 (YYYY-MM)"
            # category_axis.tick_labels.font.size = Pt(14)

            value_axis = chart.value_axis
            value_axis.has_title = True
            value_axis.axis_title.text_frame.text = "티켓 건수"
            value_axis.tick_labels.font.size = Pt(8)
            # value_axis.tick_labels.font.bold = True
            # value_axis.major_unit = 1  # 1씩 증가 (정수 단위로 눈금 표시)
            value_axis.minor_unit = 1  # (선택) 보조 눈금도 1로 설정 가능
            value_axis.has_major_gridlines = False

            

            
        elif i == 3 :
            content_slide_layout = prs.slide_layouts[13]
            slide = prs.slides.add_slide(content_slide_layout)

            title = slide.placeholders[0]
            title_tf = title.text_frame
            title_tf.text = f"2. Contacted service graph"
            set_font(title_tf.paragraphs[0], size=18, bold=True)

            # 그래프 생성
            service_names = list(service_name_counts.keys())
            service_counts = list(service_name_counts.values())

            # pptx 내장 파이 차트 데이터 생성
            chart_data_service = ChartData()
            chart_data_service.categories = service_names
            chart_data_service.add_series('서비스별 문의 비율', service_counts)

            x, y, cx, cy = Inches(0.5), Inches(1), Inches(4.5), Inches(4.2)
            service_chart = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data_service
            ).chart
            
            service_chart.has_legend = True
            service_chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # 하단(BOTTOM)
            service_chart.legend.font.size = Pt(10)  # 폰트 크기 설정
            service_chart.plots[0].has_data_labels = True
            data_labels = service_chart.plots[0].data_labels
            data_labels.show_percentage = False      # 퍼센트(%) 표시
            data_labels.show_value = True          # 건수(값) 숨김 (원하면 True로)
            data_labels.show_category_name = False   # 카테고리명(서비스명) 표시 (원하면 True로)
            data_labels.font.size = Pt(8)          # 폰트 크기 설정
            data_labels.font.color.rgb = RGBColor(255, 255, 255)  # 폰트 색상 설정<en

            account_ids = list(account_ids_counts.keys())
            account_ids_counts = list(account_ids_counts.values())

            chart_data_account = ChartData()
            chart_data_account.categories = account_ids
            chart_data_account.add_series('계정별 문의 건수', account_ids_counts)

            x, y, cx, cy = Inches(5), Inches(1), Inches(4.5), Inches(4.2)
            account_chart = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data_account
            ).chart
            
            account_chart.has_legend = True
            account_chart.legend.position = XL_LEGEND_POSITION.BOTTOM  # 하단(BOTTOM)
            account_chart.legend.font.size = Pt(10)  # 폰트 크기 설정
            account_chart.plots[0].has_data_labels = True
            data_labels = account_chart.plots[0].data_labels
            data_labels.show_percentage = False      # 퍼센트(%) 표시
            data_labels.show_value = True          # 건수(값) 숨김 (원하면 True로)
            data_labels.show_category_name = False   # 카테고리명(서비스명) 표시 (원하면 True로)
            data_labels.font.size = Pt(8)          # 폰트 크기 설정
            data_labels.font.color.rgb = RGBColor(255, 255, 255)  # 폰트 색상 설정<en





        elif i == 4 :
            print("본문? i==4에 내용 도달했는가")
            for account_index, account_obj in enumerate(json_data):

                for index, obj in enumerate(account_obj["tickets"]):
                    
                    print("본문? 내용 도달했는가2222")
                    account_id = account_obj['account_id']
                    ticket_id = obj['ticket_id']
                    subject = obj['subject']
                    service_name = obj['service_name']
                    content_slide_layout = prs.slide_layouts[13]
                    slide = prs.slides.add_slide(content_slide_layout)

                    title = slide.placeholders[0]
                    title_tf = title.text_frame
                    title_tf.text = f"3. Summary Report by Ticket - {ticket_id}"
                    set_font(title_tf.paragraphs[0], size=18, bold=True)

                    # ---add table to slide---
                    x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(0.5)
                    shape_information_table = slide.shapes.add_table(2, 3, x, y, cx, cy)

                    information_table = shape_information_table.table
                    information_table.rows[0].height = int(cy * (1/3))
                    information_table.rows[1].height = int(cy * (2/3))
                
                    set_cell(information_table.cell(0, 0), "Account ID", size=8,bold=True, column = 1)
                    set_cell(information_table.cell(0, 1), "Subject", size=8,bold=True, column = 1)
                    set_cell(information_table.cell(0, 2), "Service Name", size=8,bold=True, column = 1)


                    set_cell(information_table.cell(1, 0), account_id, size=12, fore_color=(232,236,229), font_color=(0,0,0))
                    set_cell(information_table.cell(1, 1), subject, size=12,fore_color=(232,236,229), font_color=(0,0,0))
                    set_cell(information_table.cell(1, 2), service_name, size=12,fore_color=(232,236,229), font_color=(0,0,0))

                    # table: pptx.table.Table 객체
                    for row in information_table.rows:
                        for cell in row.cells:
                            set_cell_border(cell, border_color="000000", border_width='12700')  # 검정, 1pt


                    # ---add table to slide---
                    x, y, cx, cy = Inches(0.5), Inches(1.8), Inches(9), Inches(1.3)
                    shape_request_table = slide.shapes.add_table(2, 1, x, y, cx, cy)
                    request_table = shape_request_table.table
                    request_table.rows[0].height = int(cy * (1/4))
                    request_table.rows[1].height = int(cy * (3/4))
                    set_cell(request_table.cell(0, 0), "Service Request Summary", size=14,bold=True, column = 1)

                    request_txt =""
                    for idx, req in enumerate(obj["request_summary"]):
                        if idx > 0 :
                            request_txt += "\n"
                        request_txt += f" - {req} "
                    
                    set_cell(request_table.cell(1, 0), request_txt, size=12, fore_color=(232,236,229), font_color=(0,0,0))


                    # ---add table to slide---
                    x, y, cx, cy = Inches(0.5), Inches(3.5), Inches(9), Inches(1.5)
                    shape_response_table = slide.shapes.add_table(2, 1, x, y, cx, cy)
                    response_table = shape_response_table.table
                    response_table.rows[0].height = int(cy * (1/4))
                    response_table.rows[1].height = int(cy * (3/4))
                    set_cell(response_table.cell(0, 0), "Service Response Summary", size=14,bold=True, column = 1)

                    response_txt = ""
                    for idx, res in enumerate(obj["response_summary"]):
                        if idx > 0 :
                            response_txt += "\n"
                        response_txt += f" - {res} "

                    set_cell(response_table.cell(1, 0), response_txt, size=12, fore_color=(232,236,229), font_color=(0,0,0))
                    
                    
    file_name = f"{company}_aws보고서-chart-test.pptx"
    file_path = "tmp/"+file_name

    result_stream = io.BytesIO()
    prs.save(result_stream)
    result_stream.seek(0)
    s3.upload_fileobj(result_stream, PPTX_S3_BUCKET_NAME, file_path)
    
    # Generate signed URL for the file
    url = s3.generate_presigned_url(
        'get_object',
        Params={'Bucket': PPTX_S3_BUCKET_NAME, 'Key': file_path},
        ExpiresIn=3600
    )

    return {
        'statusCode': 200,
        'fileurl' : url
    }

